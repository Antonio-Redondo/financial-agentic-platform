# Standard library imports
import os
import sys
import threading
import time
from datetime import datetime

# Third-party imports
import numpy as np
import pandas as pd
import streamlit as st
from dotenv import load_dotenv

# Add the project root directory to the Python path
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), "../.."))
if project_root not in sys.path:
    sys.path.append(project_root)

# Local imports
from src.agents.financial_agent import FinancialAgent
from src.agents.langsmith_integration import langsmith_manager

load_dotenv()

# Enhanced page configuration (must be first Streamlit command)
st.set_page_config(
    page_title="Financial Forecast AI | Prepayment Analytics",
    page_icon="🏦",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize LangSmith integration after page config
if langsmith_manager.is_enabled:
    st.success("🔬 LangSmith observability enabled")
else:
    st.info("ℹ️ LangSmith observability disabled")

# Custom CSS for modern styling
def load_custom_css():
    st.markdown("""
    <style>
    /* Simple Typography System */
    .main {
        padding-top: 2rem;
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', 'Roboto', sans-serif;
    }
    
    /* Header styling with enhanced typography */
    .main-header {
        background: linear-gradient(90deg, #0f2027 0%, #2c5364 50%, #3b82f6 100%);
        padding: 2.5rem 2rem 2rem 2rem;
        border-radius: 16px;
        margin-bottom: 2.5rem;
        color: #fff;
        text-align: center;
        box-shadow: 0 6px 24px rgba(59,130,246,0.15);
        position: relative;
        overflow: hidden;
    }
    .main-header::after {
        content: "";
        position: absolute;
        top: 0; left: 0; right: 0; bottom: 0;
        background: url('https://img.icons8.com/ios-filled/100/ffffff/artificial-intelligence.png') no-repeat right 2rem center;
        opacity: 0.08;
        pointer-events: none;
    }
    .main-header h1 {
        margin: 0;
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
        font-size: 3rem;
        font-weight: 700;
        text-shadow: 0 2px 8px rgba(0,0,0,0.15);
        color: white;
    }
    
    .main-header p {
        margin: 1rem 0 0 0;
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
        font-size: 1.3rem;
        font-weight: 400;
        opacity: 0.95;
    }
    
    /* Metrics cards */
    .metric-card {
        background: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        border-left: 4px solid #3b82f6;
        margin-bottom: 1rem;
    }
    
    /* Chat message styling */
    .chat-message {
        padding: 1rem;
        border-radius: 10px;
        margin-bottom: 1rem;
        border-left: 4px solid #3b82f6;
    }
    
    .user-message {
        background-color: #eff6ff;
        border-left-color: #3b82f6;
    }
    
    .assistant-message {
        background-color: #f0fdf4;
        border-left-color: #10b981;
    }
    
    /* Sidebar styling */
    .sidebar .sidebar-content {
        background-color: #f8fafc;
    }
    
    /* Upload area styling */
    .upload-area {
        border: 2px dashed #cbd5e1;
        border-radius: 10px;
        padding: 2rem;
        text-align: center;
        background-color: #f8fafc;
        margin-bottom: 1rem;
    }
    
    /* Analysis results styling */
    .analysis-section {
        background: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        margin-bottom: 1rem;
    }
    
    /* Status indicators */
    .status-success {
        color: #10b981;
        font-weight: 600;
    }
    
    .status-warning {
        color: #f59e0b;
        font-weight: 600;
    }
    
    .status-error {
        color: #ef4444;
        font-weight: 600;
    }
    
    /* Button styling */
    .stButton > button {
        width: 100%;
        border-radius: 10px;
        border: none;
        background: linear-gradient(90deg, #3b82f6, #1d4ed8);
        color: white;
        font-weight: 600;
        padding: 0.75rem 1.5rem;
        transition: all 0.3s ease;
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(59, 130, 246, 0.4);
    }
    
    /* Enhanced Loading Animations */
    @keyframes spin {
        0% { transform: rotate(0deg); }
        100% { transform: rotate(360deg); }
    }
    
    @keyframes pulse {
        0% { opacity: 1; }
        50% { opacity: 0.5; }
        100% { opacity: 1; }
    }
    
    @keyframes bounce {
        0%, 20%, 50%, 80%, 100% { transform: translateY(0); }
        40% { transform: translateY(-10px); }
        60% { transform: translateY(-5px); }
    }
    
    @keyframes fadeInOut {
        0% { opacity: 0; transform: scale(0.8); }
        50% { opacity: 1; transform: scale(1.1); }
        100% { opacity: 1; transform: scale(1); }
    }
    
    /* Enhanced Custom Spinner Styles */
    .custom-spinner {
        display: inline-block;
        width: 40px;
        height: 40px;
        border: 4px solid rgba(255, 255, 255, 0.3);
        border-top: 4px solid #ffffff;
        border-radius: 50%;
        animation: spinGlow 1.2s linear infinite;
        margin: 8px;
        position: relative;
        box-shadow: 0 0 20px rgba(59, 130, 246, 0.3);
    }
    
    .custom-spinner::before {
        content: '';
        position: absolute;
        top: -4px;
        left: -4px;
        right: -4px;
        bottom: -4px;
        border-radius: 50%;
        border: 2px solid transparent;
        border-top: 2px solid rgba(255, 255, 255, 0.5);
        animation: spinGlow 2s linear infinite reverse;
    }
    
    @keyframes spinGlow {
        0% { transform: rotate(0deg); box-shadow: 0 0 20px rgba(59, 130, 246, 0.3); }
        25% { box-shadow: 0 0 30px rgba(139, 92, 246, 0.4); }
        50% { transform: rotate(180deg); box-shadow: 0 0 25px rgba(236, 72, 153, 0.4); }
        75% { box-shadow: 0 0 35px rgba(59, 130, 246, 0.5); }
        100% { transform: rotate(360deg); box-shadow: 0 0 20px rgba(59, 130, 246, 0.3); }
    }
    
    .analyzing-container {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 50%, #ff6b6b 100%);
        background-size: 200% 200%;
        animation: gradientShift 3s ease infinite, slideInScale 0.8s ease-out;
        color: white;
        padding: 1.5rem 1rem;
        border-radius: 16px;
        text-align: center;
        margin: 0.5rem 0;
        box-shadow: 0 8px 32px rgba(0,0,0,0.15), 0 0 0 1px rgba(255,255,255,0.1);
        position: relative;
        overflow: hidden;
        backdrop-filter: blur(10px);
    }
    
    .analyzing-container::before {
        content: '';
        position: absolute;
        top: 0;
        left: -100%;
        width: 100%;
        height: 100%;
        background: linear-gradient(90deg, transparent, rgba(255,255,255,0.2), transparent);
        animation: shimmer 2s infinite;
    }
    
    @keyframes gradientShift {
        0% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
        100% { background-position: 0% 50%; }
    }
    
    @keyframes shimmer {
        0% { left: -100%; }
        100% { left: 100%; }
    }
    
    @keyframes slideInScale {
        0% { 
            opacity: 0; 
            transform: translateY(20px) scale(0.95); 
        }
        100% { 
            opacity: 1; 
            transform: translateY(0) scale(1); 
        }
    }
    
    .analyzing-text {
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
        font-size: 1.1rem;
        font-weight: 600;
        margin: 0.5rem 0;
        animation: textGlow 2s ease-in-out infinite alternate;
        text-shadow: 0 2px 8px rgba(0,0,0,0.2);
    }
    
    @keyframes textGlow {
        0% { 
            text-shadow: 0 2px 4px rgba(0,0,0,0.3), 0 0 10px rgba(255,255,255,0.2); 
        }
        100% { 
            text-shadow: 0 2px 4px rgba(0,0,0,0.3), 0 0 20px rgba(255,255,255,0.4), 0 0 30px rgba(59,130,246,0.3); 
        }
    }
    
    .analyzing-subtext {
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
        font-size: 0.9rem;
        font-weight: 400;
        opacity: 0.9;
        margin-top: 0.5rem;
        animation: subtlePulse 3s ease-in-out infinite;
    }
    
    @keyframes subtlePulse {
        0%, 100% { opacity: 0.95; }
        50% { opacity: 0.8; }
    }
    
    .brain-icon {
        font-size: 2.5rem;
        animation: brainPulse 2.5s ease-in-out infinite;
        margin-bottom: 0.5rem;
        display: inline-block;
        filter: drop-shadow(0 4px 8px rgba(0,0,0,0.2));
    }
    
    @keyframes brainPulse {
        0%, 100% { 
            transform: scale(1) rotate(0deg); 
            filter: drop-shadow(0 4px 8px rgba(0,0,0,0.2)); 
        }
        25% { 
            transform: scale(1.1) rotate(2deg); 
            filter: drop-shadow(0 6px 12px rgba(139,92,246,0.3)); 
        }
        50% { 
            transform: scale(1.05) rotate(0deg); 
            filter: drop-shadow(0 8px 16px rgba(59,130,246,0.4)); 
        }
        75% { 
            transform: scale(1.1) rotate(-2deg); 
            filter: drop-shadow(0 6px 12px rgba(236,72,153,0.3)); 
        }
    }
    
    /* Enhanced Progress Dots Animation */
    .progress-dots {
        display: inline-block;
        margin-left: 10px;
        position: relative;
    }
    
    .progress-dots::after {
        content: '';
        animation: enhancedDots 2s infinite;
        font-weight: bold;
        color: rgba(255,255,255,0.9);
    }
    
    @keyframes enhancedDots {
        0% { content: ''; }
        20% { content: '●'; }
        40% { content: '●●'; }
        60% { content: '●●●'; }
        80% { content: '●●●●'; }
        100% { content: '●●●●●'; }
    }
    
    /* Search Processing Background Effect */
    .search-processing-bg {
        position: fixed;
        top: 0;
        left: 0;
        width: 100%;
        height: 100%;
        background: linear-gradient(45deg, rgba(59,130,246,0.05), rgba(139,92,246,0.05), rgba(236,72,153,0.05));
        background-size: 400% 400%;
        animation: searchBgShift 4s ease infinite;
        pointer-events: none;
        z-index: -1;
    }
    
    @keyframes searchBgShift {
        0% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
        100% { background-position: 0% 50%; }
    }
    
    /* Enhanced Success Animation */
    .success-animation {
        background: linear-gradient(135deg, #10b981 0%, #059669 100%);
        color: white;
        padding: 1rem;
        border-radius: 12px;
        text-align: center;
        margin: 0.5rem 0;
        box-shadow: 0 4px 20px rgba(16, 185, 129, 0.3);
        animation: successPulse 0.6s ease-out;
        position: relative;
        overflow: hidden;
    }
    
    .success-animation::before {
        content: '';
        position: absolute;
        top: 0;
        left: -100%;
        width: 100%;
        height: 100%;
        background: linear-gradient(90deg, transparent, rgba(255,255,255,0.3), transparent);
        animation: successShimmer 1s ease-out;
    }
    
    @keyframes successPulse {
        0% { transform: scale(0.95); opacity: 0; }
        50% { transform: scale(1.02); }
        100% { transform: scale(1); opacity: 1; }
    }
    
    @keyframes successShimmer {
        0% { left: -100%; }
        100% { left: 100%; }
    }
    
    /* Typing Effect for Search Results */
    .typing-effect {
        overflow: hidden;
        border-right: 2px solid #3b82f6;
        animation: typewriter 2s steps(40) 1s both, blinkCursor 1s infinite;
        white-space: nowrap;
    }
    
    @keyframes typewriter {
        from { width: 0; }
        to { width: 100%; }
    }
    
    @keyframes blinkCursor {
        50% { border-color: transparent; }
    }
        50% { content: '..'; }
        75% { content: '...'; }
        100% { content: ''; }
    }
    
    /* Success Animation */
    .success-animation {
        animation: fadeInOut 0.8s ease-out;
        background: linear-gradient(135deg, #10b981 0%, #059669 100%);
        color: white;
        padding: 1rem;
        border-radius: 10px;
        text-align: center;
        margin: 1rem 0;
    }
    
    /* Enhanced Chat Styling with Simple Typography */
    .chat-message {
        padding: 1.8rem;
        border-radius: 16px;
        margin-bottom: 1.5rem;
        box-shadow: 0 3px 12px rgba(0,0,0,0.08);
        animation: fadeInOut 0.5s ease-out;
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
        line-height: 1.6;
    }
    
    .user-message {
        background: linear-gradient(135deg, #eff6ff 0%, #dbeafe 100%);
        border-left: 5px solid #3b82f6;
        font-weight: 500;
    }
    
    .assistant-message {
        background: linear-gradient(135deg, #f0fdf4 0%, #dcfce7 100%);
        border-left: 5px solid #10b981;
        font-weight: 400;
    }
    
    /* Simple Typography for Messages */
    .chat-message p {
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
        font-size: 1rem;
        line-height: 1.6;
        margin-bottom: 0.8rem;
        color: #374151;
    }
    
    .chat-message h1, .chat-message h2, .chat-message h3 {
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
        font-weight: 600;
        margin-top: 1.5rem;
        margin-bottom: 0.8rem;
    }
    
    .chat-message h1 { font-size: 1.8rem; color: #1f2937; }
    .chat-message h2 { font-size: 1.5rem; color: #374151; }
    .chat-message h3 { font-size: 1.3rem; color: #4b5563; }
    
    .chat-message code {
        font-family: 'JetBrains Mono', 'SF Mono', 'Monaco', 'Consolas', monospace;
        background: rgba(99, 102, 241, 0.1);
        padding: 0.2rem 0.4rem;
        border-radius: 4px;
        font-size: 0.9rem;
        font-weight: 500;
        letter-spacing: 0;
    }
    
    .chat-message pre {
        font-family: 'JetBrains Mono', 'SF Mono', 'Monaco', 'Consolas', monospace;
        background: #f8fafc;
        padding: 1rem;
        border-radius: 8px;
        border-left: 4px solid #e2e8f0;
        overflow-x: auto;
        font-size: 0.9rem;
        line-height: 1.5;
    }
    
    /* Typing Animation */
    .typing-indicator {
        display: inline-block;
        background: #e5e7eb;
        padding: 10px 15px;
        border-radius: 15px;
        margin: 10px 0;
    }
    
    .typing-indicator::after {
        content: '';
        display: inline-block;
        width: 3px;
        height: 3px;
        border-radius: 50%;
        background: #6b7280;
        animation: typing 1.4s infinite;
        margin-left: 5px;
    }
    
    @keyframes typing {
        0% { box-shadow: 10px 0 #6b7280, 20px 0 #6b7280; }
        25% { box-shadow: 10px 0 #3b82f6, 20px 0 #6b7280; }
        50% { box-shadow: 10px 0 #6b7280, 20px 0 #3b82f6; }
        75% { box-shadow: 10px 0 #6b7280, 20px 0 #6b7280; }
        100% { box-shadow: 10px 0 #6b7280, 20px 0 #6b7280; }
    }
    
    /* Hide Streamlit branding */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    
    /* Enhanced Chat Input Styling with Typography */
    .stChatInput > div > div > div > div {
        min-height: 85px !important;
        height: 85px !important;
        background: transparent !important;
        background-color: transparent !important;
    }
    
    .stChatInput textarea {
        min-height: 65px !important;
        height: 65px !important;
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif !important;
        font-size: 16px !important;
        font-weight: 400 !important;
        padding: 16px 20px !important;
        line-height: 1.5 !important;
        resize: vertical !important;
        background: transparent !important;
        background-color: transparent !important;
        color: #374151 !important;
        transition: all 0.3s ease !important;
    }
    
    .stChatInput > div {
        background: transparent !important;
        background-color: transparent !important;
        border: 2px solid #e2e8f0 !important;
        border-radius: 12px !important;
        box-shadow: 0 2px 8px rgba(0,0,0,0.1) !important;
        margin: 8px 0 !important;
    }
    
    .stChatInput > div > div {
        background: transparent !important;
        background-color: transparent !important;
    }
    
    .stChatInput > div > div > div {
        background: transparent !important;
        background-color: transparent !important;
    }
    
    .stChatInput > div > div > div > div {
        background: transparent !important;
        background-color: transparent !important;
    }
    
    .stChatInput > div:focus-within {
        border-color: #3b82f6 !important;
        box-shadow: 0 0 0 3px rgba(59, 130, 246, 0.1) !important;
    }
    
    /* Simple Typography and Font Effects */
    .stChatInput textarea::placeholder {
        color: #9ca3af !important;
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif !important;
        font-size: 15px !important;
        font-weight: 400 !important;
        opacity: 0.75 !important;
        transition: all 0.3s ease !important;
    }
    
    .stChatInput textarea:focus::placeholder {
        opacity: 0.4 !important;
        color: #6b7280 !important;
    }
    
    .stChatInput textarea:focus {
        font-weight: 500 !important;
        color: #1f2937 !important;
    }
    
    /* Search Result Highlighting */
    .search-highlight {
        background: linear-gradient(120deg, rgba(59,130,246,0.1) 0%, rgba(139,92,246,0.1) 100%);
        padding: 2px 4px;
        border-radius: 4px;
        font-weight: 600;
        animation: highlightPulse 2s ease-in-out infinite;
    }
    
    @keyframes highlightPulse {
        0%, 100% { background: linear-gradient(120deg, rgba(59,130,246,0.1) 0%, rgba(139,92,246,0.1) 100%); }
        50% { background: linear-gradient(120deg, rgba(59,130,246,0.2) 0%, rgba(139,92,246,0.2) 100%); }
    }
    </style>
    """, unsafe_allow_html=True)

def create_header():
    st.markdown("""
    <div class="main-header">
        <h1>🤖 AI Financial Analyst</h1>
        <p>Advanced Prepayment Analytics, Bond Valuation & Risk Assessment</p>
        <div style="margin-top: 1.2rem; font-size: 0.95rem; opacity: 0.9; font-weight: 400;">
            <span style="background: rgba(255,255,255,0.2); padding: 0.4rem 0.9rem; border-radius: 25px; margin: 0 0.3rem; font-weight: 500;">📊 Real-time Analysis</span>
            <span style="background: rgba(255,255,255,0.2); padding: 0.4rem 0.9rem; border-radius: 25px; margin: 0 0.3rem; font-weight: 500;">🧠 AI-Powered</span>
            <span style="background: rgba(255,255,255,0.2); padding: 0.4rem 0.9rem; border-radius: 25px; margin: 0 0.3rem; font-weight: 500;">⚡ Lightning Fast</span>
        </div>
    </div>
    """, unsafe_allow_html=True)

def create_metrics_dashboard():
    col1, col2 = st.columns(2)
    
    with col1:
        st.metric(
            label=" Messages",
            value=len(st.session_state.get("messages", [])),
            delta=f"+{len(st.session_state.get('messages', [])) // 2} conversations" if len(st.session_state.get("messages", [])) > 0 else None
        )
    
    with col2:
        agent_status = "🟢 Ready" if st.session_state.get("agent") else "🔴 Error"
        st.metric(
            label="🎯 System Status",
            value="Operational" if st.session_state.get("agent") else "Offline",
            delta=agent_status
        )

def create_sidebar():
    with st.sidebar:
        st.markdown("### 📁 Document Management")
        
        # Upload section with better styling
        st.markdown('<div class="upload-area">', unsafe_allow_html=True)
        uploaded_files = st.file_uploader(
            "📄 Upload Financial Documents",
            accept_multiple_files=True,
            type=["pdf", "docx", "doc", "txt", "rtf", "csv", "xlsx", "xls", "pptx", "ppt", "html", "htm", "xml", "json", "md"],
            help="Supported formats: PDF, Word (DOCX/DOC), Excel (XLSX/XLS), PowerPoint (PPTX/PPT), Text (TXT/RTF/MD), Web (HTML/XML), Data (CSV/JSON). Multiple files supported."
        )
        st.markdown('</div>', unsafe_allow_html=True)
        
        if uploaded_files:
            st.markdown("#### 📋 Uploaded Files")
            for i, file in enumerate(uploaded_files):
                file_size = len(file.getvalue()) / 1024  # KB
                file_extension = file.name.split('.')[-1].upper()
                
                # Choose icon based on file type
                if file_extension in ['PDF']:
                    icon = "📕"
                elif file_extension in ['DOCX', 'DOC']:
                    icon = "📘"
                elif file_extension in ['XLSX', 'XLS', 'CSV']:
                    icon = "📊"
                elif file_extension in ['PPTX', 'PPT']:
                    icon = "📽️"
                elif file_extension in ['TXT', 'MD', 'RTF']:
                    icon = "📝"
                elif file_extension in ['HTML', 'HTM', 'XML']:
                    icon = "🌐"
                elif file_extension in ['JSON']:
                    icon = "📋"
                else:
                    icon = "📄"
                
                st.markdown(f"""
                <div style="background: #f0fdf4; padding: 0.5rem; border-radius: 5px; margin: 0.25rem 0;">
                    {icon} {file.name}<br>
                    <small style="color: #6b7280;">Type: {file_extension} | Size: {file_size:.1f} KB</small>
                </div>
                """, unsafe_allow_html=True)
            
            if st.button("🚀 Process Documents", type="primary"):
                process_documents(uploaded_files)
        
        # Show processed documents status
        if st.session_state.get("processed_docs"):
            st.markdown("#### 📊 Processed Documents")
            for doc in st.session_state.processed_docs:
                icon = "✅" if doc.get("indexed", False) else "⚠️"
                chunks = doc.get("chunks_created", 0)
                file_ext = doc['name'].split('.')[-1].upper()
                
                st.markdown(f"""
                <div style="background: #f8fafc; padding: 0.5rem; border-radius: 5px; margin: 0.25rem 0; border-left: 3px solid {'#10b981' if doc.get('indexed') else '#f59e0b'};">
                    {icon} {doc['name']}<br>
                    <small style="color: #6b7280;">
                        Type: {file_ext} | 
                        {'Indexed: ' + str(chunks) + ' chunks' if doc.get('indexed') else 'Not indexed'} | 
                        Size: {doc['size'] / 1024:.1f} KB | 
                        Content: {len(doc.get('content', '')):.0f} chars
                    </small>
                </div>
                """, unsafe_allow_html=True)
            
            # Add document search test
            if st.button("🔍 Test Document Search", help="Test if documents are searchable"):
                if st.session_state.agent and st.session_state.agent.vector_store:
                    try:
                        stats = st.session_state.agent.get_document_stats()
                        st.success(f"📊 Document Stats: {stats['total_chunks']} chunks from {stats['unique_documents']} documents in {stats['storage_type']} storage")
                        
                        # Show all documents info for debugging
                        if hasattr(st.session_state.agent.vector_store, 'get_all_documents_info'):
                            all_docs = st.session_state.agent.vector_store.get_all_documents_info()
                            if all_docs:
                                with st.expander("📋 All Stored Documents (Debug Info)"):
                                    for i, doc_info in enumerate(all_docs, 1):
                                        st.text(f"Document {i}:")
                                        st.text(f"  Filename: {doc_info.get('filename', 'unknown')}")
                                        st.text(f"  Content Length: {doc_info.get('content_length', 0)} chars")
                                        st.text(f"  Preview: {doc_info.get('content_preview', 'No preview')}")
                                        st.text("---")
                        
                        # Test search with user input
                        test_query = st.text_input("🔍 Test search query:", value="financial", help="Enter a word to search for in your documents")
                        if st.button("Search Documents"):
                            test_results = st.session_state.agent.vector_store.search_documents(test_query, k=5)
                            if test_results:
                                st.success(f"🔍 Search Test: Found {len(test_results)} relevant chunks for '{test_query}'")
                                with st.expander("Preview search results"):
                                    for i, result in enumerate(test_results, 1):
                                        st.text(f"Result {i} (Relevance: {result.get('similarity_score', 0):.2f}):")
                                        st.text(f"Source: {result.get('metadata', {}).get('filename', 'unknown')}")
                                        st.text(f"Content: {result.get('content', '')[:300]}...")
                                        st.text("---")
                            else:
                                st.warning(f"🔍 Search Test: No results found for '{test_query}' - documents may not contain this content")
                        
                    except Exception as e:
                        st.error(f"❌ Error testing document search: {str(e)}")
                else:
                    st.error("❌ Vector store not available")
        
        st.markdown("---")
        st.markdown("### 📊 Document Stats")
        st.markdown("""
        - **Total chunks:** 387
        - **Deals covered:** 3
        - **Database:** PostgreSQL + pgvector
        """)
        
        # S3 Knowledge Base Status
        st.markdown("### 🗂️ Knowledge Base Status")
        
        # Check for background processing completion
        check_s3_background_status()
        
        # Show S3 status based on current state
        s3_status = st.session_state.get("s3_status", "ready")
        
        if s3_status == "ready":
            st.info("📋 Ready to sync with S3")
            if st.button("🚀 Start S3 Sync", type="primary"):
                start_s3_background_processing()
                st.rerun()
        elif s3_status == "processing":
            st.info("🔄 Loading documents from S3...")
            # Auto-refresh every 3 seconds while processing
            time.sleep(3)
            st.rerun()
        elif s3_status == "completed":
            results = st.session_state.get("s3_init_results", {})
            if results.get("documents_processed", 0) > 0:
                st.success(f"✅ S3 Auto-load: {results['documents_processed']} documents indexed")
                
                with st.expander("� View S3 Document Details"):
                    st.write(f"**Processing Time:** {results.get('processing_time', 0):.2f}s")
                    st.write(f"**Total Documents Found:** {results.get('total_documents', 0)}")
                    st.write(f"**New Documents Indexed:** {results['documents_processed']}")
                    st.write(f"**Duplicates Skipped:** {results.get('skipped_duplicates', 0)}")
                    
                    if results.get('processed_documents'):
                        st.write("**Processed Documents:**")
                        for doc in results['processed_documents']:
                            st.text(f"• {doc['filename']} ({doc['chunks']} chunks, {doc['size']} chars)")
            else:
                st.info("ℹ️ S3 bucket is empty or no new documents found")
                
            if st.button("🔄 Refresh from S3"):
                # Reset status to trigger re-processing
                st.session_state.s3_status = "ready"
                st.session_state.s3_background_started = False
                st.rerun()
                
        elif s3_status == "error":
            error_msg = st.session_state.get("s3_init_error", "Unknown error")
            st.error(f"❌ S3 processing failed: {error_msg}")
            
            if st.button("🔄 Retry S3 Processing"):
                st.session_state.s3_status = "ready"
                st.session_state.s3_background_started = False
                st.session_state.s3_init_error = None
                st.rerun()
                
        elif s3_status == "disabled":
            st.warning("⚠️ S3 auto-loading is disabled")
            st.write("To enable S3 auto-loading:")
            st.code("S3_AUTO_LOAD_ENABLED=true")
        else:
            st.warning(f"⚠️ Unknown S3 status: {s3_status}")
        
        # LangSmith Analytics Section
        if langsmith_manager.is_enabled:
            st.markdown("---")
            st.markdown("### 🔬 LangSmith Analytics")
            
            # Get project analytics
            try:
                analytics = langsmith_manager.get_project_analytics(days_back=7)
                
                if analytics and analytics.get("status") == "full_access":
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        st.metric("Total Runs", analytics.get("total_runs", 0))
                        st.metric("Success Rate", f"{(analytics.get('successful_runs', 0) / max(analytics.get('total_runs', 1), 1) * 100):.1f}%")
                    
                    with col2:
                        st.metric("Error Runs", analytics.get("error_runs", 0))
                        avg_latency = analytics.get("avg_latency", 0)
                        st.metric("Avg Latency", f"{avg_latency:.2f}s" if avg_latency else "N/A")
                    
                    st.markdown(f"**Project:** {analytics.get('project_name', 'Unknown')}")
                    st.markdown(f"**Last 7 days**")
                    
                    # Link to LangSmith dashboard
                    langsmith_url = f"https://smith.langchain.com/o/{os.getenv('LANGCHAIN_PROJECT', 'financial-forecast-ai-app')}"
                    st.markdown(f"🔗 [View Full Dashboard]({langsmith_url})")
                    
                elif analytics and analytics.get("status") in ["limited_access", "access_denied"]:
                    st.info("📊 LangSmith Analytics")
                    st.markdown(f"**Status:** {analytics.get('message', 'Limited access')}")
                    st.markdown(f"**Local Tracing:** {analytics.get('local_tracing', 'enabled')}")
                    st.markdown("**Features Available:**")
                    st.markdown("- ✅ Query and response tracking")
                    st.markdown("- ✅ Performance monitoring") 
                    st.markdown("- ✅ Error logging")
                    st.markdown("- ⚠️ Dashboard analytics (limited)")
                    
                    # Link to LangSmith dashboard
                    langsmith_url = f"https://smith.langchain.com/o/{os.getenv('LANGCHAIN_PROJECT', 'financial-forecast-ai-app')}"
                    st.markdown(f"🔗 [View LangSmith Dashboard]({langsmith_url})")
                else:
                    st.info("📊 No analytics data available yet")
                    
            except Exception as e:
                st.warning(f"⚠️ LangSmith analytics unavailable: {str(e)}")
        else:
            st.markdown("---")
            st.markdown("### 🔬 LangSmith Integration")
            st.info("💡 Enable LangSmith tracing by setting LANGCHAIN_TRACING_V2=true in your .env file")
        
        st.markdown("---")

def process_documents(uploaded_files):
    """Process uploaded documents and save to vector database"""
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    # Get the vector store from the agent
    vector_store = st.session_state.agent.vector_store if st.session_state.agent else None
    
    for i, file in enumerate(uploaded_files):
        status_text.text(f"Processing {file.name}...")
        
        # Read file content based on file type
        try:
            content = ""
            file_extension = file.name.split('.')[-1].lower()
            
            if file.type == "application/pdf" or file_extension == "pdf":
                # PDF processing with pypdf
                try:
                    import pypdf
                    import io
                    
                    pdf_reader = pypdf.PdfReader(io.BytesIO(file.getvalue()))
                    for page in pdf_reader.pages:
                        content += page.extract_text() + "\n"
                    
                    if not content.strip():
                        content = "PDF content could not be extracted (possibly scanned document)"
                except Exception as e:
                    content = f"Error processing PDF: {str(e)}"
                    
            elif file.type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document"] or file_extension == "docx":
                # DOCX files
                try:
                    from docx import Document
                    import io
                    
                    doc = Document(io.BytesIO(file.getvalue()))
                    for paragraph in doc.paragraphs:
                        content += paragraph.text + "\n"
                    
                    # Also extract text from tables
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                content += cell.text + " "
                            content += "\n"
                    
                    if not content.strip():
                        content = "DOCX content could not be extracted"
                except Exception as e:
                    content = f"Error processing DOCX: {str(e)}"
                    
            elif file.type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"] or file_extension in ["xlsx", "xls"]:
                # Excel files
                try:
                    import pandas as pd
                    import io
                    
                    # Read all sheets
                    df_dict = pd.read_excel(io.BytesIO(file.getvalue()), sheet_name=None)
                    for sheet_name, df in df_dict.items():
                        content += f"Sheet: {sheet_name}\n"
                        content += df.to_string(index=False) + "\n\n"
                        
                except Exception as e:
                    content = f"Error processing Excel file: {str(e)}"
                    
            elif file.type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation"] or file_extension in ["pptx", "ppt"]:
                # PowerPoint files
                try:
                    from pptx import Presentation
                    import io
                    
                    prs = Presentation(io.BytesIO(file.getvalue()))
                    for slide_num, slide in enumerate(prs.slides, 1):
                        content += f"Slide {slide_num}:\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                        content += "\n"
                        
                except Exception as e:
                    content = f"Error processing PowerPoint file: {str(e)}"
                    
            elif file.type == "text/csv" or file_extension == "csv":
                # CSV files
                try:
                    import pandas as pd
                    import io
                    
                    df = pd.read_csv(io.BytesIO(file.getvalue()))
                    content = f"CSV Data ({len(df)} rows):\n"
                    content += df.to_string(index=False)
                    
                except Exception as e:
                    content = f"Error processing CSV: {str(e)}"
                    
            elif file.type == "application/json" or file_extension == "json":
                # JSON files
                try:
                    import json
                    
                    json_data = json.loads(file.getvalue().decode('utf-8'))
                    content = "JSON Content:\n" + json.dumps(json_data, indent=2)
                    
                except Exception as e:
                    content = f"Error processing JSON: {str(e)}"
                    
            elif file.type in ["text/html", "application/xml"] or file_extension in ["html", "htm", "xml"]:
                # HTML/XML files
                try:
                    from bs4 import BeautifulSoup
                    
                    html_content = file.getvalue().decode('utf-8')
                    soup = BeautifulSoup(html_content, 'html.parser')
                    content = soup.get_text()
                    
                except Exception as e:
                    content = f"Error processing HTML/XML: {str(e)}"
                    
            elif file_extension == "md":
                # Markdown files
                try:
                    import markdown
                    
                    md_content = file.getvalue().decode('utf-8')
                    html = markdown.markdown(md_content)
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(html, 'html.parser')
                    content = soup.get_text()
                    
                except Exception as e:
                    content = f"Error processing Markdown: {str(e)}"
                    
            elif file_extension == "rtf":
                # RTF files (basic extraction)
                try:
                    rtf_content = file.getvalue().decode('utf-8', errors='ignore')
                    # Basic RTF text extraction (removes most RTF codes)
                    import re
                    content = re.sub(r'\\[a-z]+\d*\s?', ' ', rtf_content)
                    content = re.sub(r'[{}]', '', content)
                    content = re.sub(r'\s+', ' ', content).strip()
                    
                except Exception as e:
                    content = f"Error processing RTF: {str(e)}"
                    
            else:
                # Default: treat as text file
                try:
                    content = file.getvalue().decode('utf-8')
                except UnicodeDecodeError:
                    try:
                        content = file.getvalue().decode('latin-1')
                    except Exception as e:
                        content = f"Could not decode file content: {str(e)}"
            
            # Create metadata for the document
            file_extension = file.name.split('.')[-1].lower()
            metadata = {
                "filename": file.name,
                "file_type": file.type,
                "file_extension": file_extension,
                "size": len(file.getvalue()),
                "uploaded_at": datetime.now().isoformat(),
                "source": "user_upload",
                "content_length": len(content)
            }
            
            # Index document into vector store (chunks automatically)
            if vector_store:
                status_text.text(f"Indexing {file.name} into vector database...")
                vector_store.index_document(content, metadata)
                status_msg = f"✅ {file.name} indexed into vector database"
            else:
                status_msg = f"⚠️ {file.name} processed but vector store unavailable"
            
            # Add to processed documents session state (with summary)
            if "processed_docs" not in st.session_state:
                st.session_state.processed_docs = []
            
            st.session_state.processed_docs.append({
                "name": file.name,
                "size": len(file.getvalue()),
                "processed_at": datetime.now(),
                "content": content[:500] + "..." if len(content) > 500 else content,
                "indexed": vector_store is not None,
                "chunks_created": len(vector_store.smart_splitter.split_text(content)) if vector_store else 0
            })
            
            # Update progress
            progress_bar.progress((i + 1) / len(uploaded_files))
            st.success(status_msg)
            
        except Exception as e:
            error_msg = f"❌ Error processing {file.name}: {str(e)}"
            st.error(error_msg)
            print(error_msg)
    
    status_text.text("✅ All documents processed and indexed successfully!")
    time.sleep(2)
    status_text.empty()
    progress_bar.empty()

def create_chat_interface():
    """Enhanced chat interface with input at top"""
    st.markdown("### 💬 AI Financial Analyst")
    
    # Create a container for the input that stays at the top
    input_container = st.container()
    
    with input_container:
        st.markdown("#### ✍️ Ask Your Question")
        
        # Chat input at the top - always visible
        prompt = st.chat_input("", key="chat_input_top")
    
    # Separator
    st.markdown("---")
    
    # Process input immediately if provided
    if prompt:
        # Add user message to chat history
        st.session_state.messages.append({"role": "user", "content": prompt})
        
        # Sync conversation with agent's LangChain memory
        if st.session_state.agent:
            try:
                # Add user message to agent's memory systems
                from langchain_core.messages import HumanMessage
                user_message = HumanMessage(content=prompt)
                
                if st.session_state.agent.buffer_memory:
                    st.session_state.agent.buffer_memory.chat_memory.add_user_message(user_message)
                    
                if st.session_state.agent.summary_memory:
                    st.session_state.agent.summary_memory.chat_memory.add_user_message(user_message)
                    
                print(f"🧠 Added user message to LangChain memory: {prompt[:50]}...")
            except Exception as e:
                print(f"⚠️ Failed to sync message to LangChain memory: {e}")
        
        # Create analysis container with fancy animations
        analysis_container = st.empty()
        
        # Show enhanced analyzing status with animations and background effect
        analysis_container.markdown("""
        <div class="search-processing-bg"></div>
        <div class="analyzing-container">
            <div class="brain-icon">🧠</div>
            <div class="analyzing-text">Analyzing your query<span class="progress-dots"></span></div>
            <div class="analyzing-subtext">🔍 AI is processing your financial request with smart document analysis</div>
            <div class="custom-spinner"></div>
            <div style="margin-top: 0.5rem; font-size: 0.7rem; opacity: 0.8;">
                ⚡ Enhanced search • 🎯 Smart chunking • 🚀 Vector analysis
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        # Process the query and get response
        import time
        try:
            start_time = time.time()
            
            # Get response from agent
            response = st.session_state.agent.process_query(prompt)
            
            # Ensure spinner is visible for at least 2 seconds for better UX
            elapsed = time.time() - start_time
            if elapsed < 2.0:
                time.sleep(2.0 - elapsed)

            # Show enhanced success animation briefly
            analysis_container.markdown("""
            <div class="success-animation">
                <div style="font-size: 2.5rem; margin-bottom: 0.5rem;">✨</div>
                <div style="font-size: 1.3rem; font-weight: bold; letter-spacing: 0.5px;">Analysis Complete!</div>
                <div style="font-size: 0.9rem; margin-top: 0.3rem;">🎯 Your financial insights are ready</div>
                <div style="font-size: 0.75rem; margin-top: 0.2rem; opacity: 0.9;">⚡ Powered by enhanced AI search</div>
            </div>
            """, unsafe_allow_html=True)
            
            # Brief pause to show success
            time.sleep(1.0)
            
            # Clear the analysis container
            analysis_container.empty()

            # Format response with simple completion message
            formatted_response = f"""**Analysis Complete** ✅

{response.get("output", "No response generated")}

*Analysis completed at {datetime.now().strftime("%H:%M:%S")}*"""

            # Add assistant message to chat history
            st.session_state.messages.append({"role": "assistant", "content": formatted_response})
            
            # Sync assistant response with agent's LangChain memory
            if st.session_state.agent:
                try:
                    from langchain_core.messages import AIMessage
                    ai_message = AIMessage(content=response.get("output", "No response generated"))
                    
                    if st.session_state.agent.buffer_memory:
                        st.session_state.agent.buffer_memory.chat_memory.add_ai_message(ai_message)
                        
                    if st.session_state.agent.summary_memory:
                        st.session_state.agent.summary_memory.chat_memory.add_ai_message(ai_message)
                        
                    print(f"🧠 Added AI response to LangChain memory")
                except Exception as e:
                    print(f"⚠️ Failed to sync AI response to LangChain memory: {e}")
            
            # Add LangSmith feedback collection
            if langsmith_manager.is_enabled:
                with st.expander("📝 Rate this response (helps improve the AI)", expanded=False):
                    feedback_col1, feedback_col2 = st.columns([3, 1])
                    
                    with feedback_col1:
                        feedback_score = st.select_slider(
                            "How helpful was this response?",
                            options=[1, 2, 3, 4, 5],
                            value=3,
                            format_func=lambda x: ["⭐", "⭐⭐", "⭐⭐⭐", "⭐⭐⭐⭐", "⭐⭐⭐⭐⭐"][x-1],
                            key=f"feedback_score_{len(st.session_state.messages)}"
                        )
                        
                        feedback_text = st.text_area(
                            "Additional feedback (optional):",
                            placeholder="What could be improved? What was particularly helpful?",
                            key=f"feedback_text_{len(st.session_state.messages)}",
                            height=80
                        )
                    
                    with feedback_col2:
                        if st.button("Submit Feedback", key=f"feedback_submit_{len(st.session_state.messages)}"):
                            session_id = st.session_state.get('session_id', str(time.time()))
                            success = langsmith_manager.log_user_feedback(
                                session_id=session_id,
                                query=prompt,
                                response=response.get("output", ""),
                                feedback_score=feedback_score,
                                feedback_text=feedback_text if feedback_text.strip() else None
                            )
                            
                            if success:
                                st.success("🙏 Thank you for your feedback!")
                            else:
                                st.warning("⚠️ Feedback not logged (LangSmith unavailable)")

            # Show additional visualizations if data is present
            if "data" in response and response["data"]:
                create_data_visualization(response["data"])

        except Exception as e:
            # Clear analyzing animation and show error
            analysis_container.empty()
            error_response = f"❌ I encountered an error while processing your request: {str(e)}"
            st.session_state.messages.append({"role": "assistant", "content": error_response})

        # Rerun to show the new messages
        st.rerun()
    
    # Show conversation history below the input
    if st.session_state.messages:
        st.markdown("#### 🔄 Latest Conversation")
        
        # Show the last user message and AI response pair
        messages_to_show = st.session_state.messages[-2:] if len(st.session_state.messages) >= 2 else st.session_state.messages[-1:]
        
        for message in messages_to_show:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
        
        st.markdown("---")
        
        # Show full conversation history in an expander
        if len(st.session_state.messages) > 2:
            with st.expander(f"📜 Full Conversation History ({len(st.session_state.messages)} messages)", expanded=False):
                for message in st.session_state.messages[:-2]:
                    with st.chat_message(message["role"]):
                        st.markdown(message["content"])
        
        # Context management controls
        col1, col2, col3 = st.columns([2, 1, 1])
        with col2:
            if st.button("🔄 Clear Context", help="Reset conversation context while keeping chat history"):
                if st.session_state.get("agent") and hasattr(st.session_state.agent, 'conversation_context'):
                    st.session_state.agent.conversation_context.clear()
                    st.session_state.agent.conversation_history.clear()
                    st.success("Context cleared!")
                    st.rerun()
        
        with col3:
            if st.button("🗑️ Clear All", help="Clear both context and chat history"):
                if st.session_state.get("agent") and hasattr(st.session_state.agent, 'conversation_context'):
                    st.session_state.agent.conversation_context.clear()
                    st.session_state.agent.conversation_history.clear()
                st.session_state.messages = []
                st.success("All cleared!")
                st.rerun()
                
    else:
        # Enhanced welcome message with animations
        st.markdown("""
        <div style="text-align: center; padding: 2rem; background: linear-gradient(135deg, #f8fafc 0%, #e2e8f0 100%); border-radius: 15px; margin: 1rem 0;">
            <div style="font-size: 3rem; margin-bottom: 1rem; animation: bounce 2s infinite;">👋</div>
            <div style="font-size: 1.5rem; font-weight: bold; color: #1e293b; margin-bottom: 0.5rem;">Welcome to AI Financial Analyst!</div>
            <div style="color: #64748b; font-size: 1.1rem;">Ask me anything about financial analysis, prepayment forecasting, or risk assessment.</div>
        </div>
        """, unsafe_allow_html=True)


def create_data_visualization(data):
    """Create visualizations for analysis data using Streamlit built-in charts"""
    if isinstance(data, dict) and data:
        st.markdown("### 📊 Analysis Results")
        
        # Display the actual data provided by the AI model
        if "chart_data" in data:
            st.markdown("#### 📈 Analysis Chart")
            chart_df = pd.DataFrame(data["chart_data"])
            st.line_chart(chart_df)
        
        if "metrics" in data:
            st.markdown("#### 📋 Key Metrics")
            metrics_df = pd.DataFrame(data["metrics"])
            st.table(metrics_df)
        
        if "summary" in data:
            st.markdown("#### 📄 Summary")
            st.write(data["summary"])
    
    elif isinstance(data, list) and data:
        st.markdown("### 📊 Analysis Results")
        st.markdown("#### 📋 Analysis Data")
        df = pd.DataFrame(data)
        st.table(df)

def sync_messages_to_langchain_memory():
    """Sync existing Streamlit messages with agent's LangChain memory"""
    if not st.session_state.agent or not st.session_state.messages:
        return
    
    try:
        from langchain_core.messages import HumanMessage, AIMessage
        
        # Clear existing memory to avoid duplicates
        if st.session_state.agent.buffer_memory:
            st.session_state.agent.buffer_memory.clear()
        if st.session_state.agent.summary_memory:
            st.session_state.agent.summary_memory.clear()
        
        # Add all existing messages to LangChain memory
        for msg in st.session_state.messages:
            if msg["role"] == "user":
                user_message = HumanMessage(content=msg["content"])
                if st.session_state.agent.buffer_memory:
                    st.session_state.agent.buffer_memory.chat_memory.add_user_message(user_message)
                if st.session_state.agent.summary_memory:
                    st.session_state.agent.summary_memory.chat_memory.add_user_message(user_message)
            elif msg["role"] == "assistant":
                # Extract the actual response content (remove formatting)
                content = msg["content"]
                if "**Analysis Complete**" in content:
                    # Extract the actual response between the header and timestamp
                    lines = content.split('\n')
                    actual_content = []
                    skip_next = False
                    for line in lines:
                        if line.startswith("**Analysis Complete**") or line.startswith("*Analysis completed at"):
                            skip_next = True
                            continue
                        if not skip_next and line.strip():
                            actual_content.append(line)
                        skip_next = False
                    content = '\n'.join(actual_content).strip()
                
                ai_message = AIMessage(content=content)
                if st.session_state.agent.buffer_memory:
                    st.session_state.agent.buffer_memory.chat_memory.add_ai_message(ai_message)
                if st.session_state.agent.summary_memory:
                    st.session_state.agent.summary_memory.chat_memory.add_ai_message(ai_message)
        
        if st.session_state.messages:
            print(f"🧠 Synced {len(st.session_state.messages)} messages to LangChain memory")
            
    except Exception as e:
        print(f"⚠️ Failed to sync messages to LangChain memory: {e}")


def initialize_session_state():
    """Initialize session state variables"""
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "agent" not in st.session_state:
        try:
            st.session_state.agent = FinancialAgent()
            print("✅ FinancialAgent initialized successfully")
            
            # Sync existing Streamlit messages with LangChain memory
            sync_messages_to_langchain_memory()
            
        except Exception as e:
            st.error(f"Failed to initialize agent: {str(e)}")
            st.session_state.agent = None
    
    if "processed_docs" not in st.session_state:
        st.session_state.processed_docs = []
    if "conversation_count" not in st.session_state:
        st.session_state.conversation_count = 0
    
    # Initialize S3 status tracking (but don't start processing)
    if "s3_status" not in st.session_state:
        st.session_state.s3_status = "ready"  # Ready to start when requested
    if "s3_init_results" not in st.session_state:
        st.session_state.s3_init_results = None
    if "s3_init_error" not in st.session_state:
        st.session_state.s3_init_error = None
    if "s3_background_started" not in st.session_state:
        st.session_state.s3_background_started = False

def s3_background_thread(agent, callback):
    """Background thread function for S3 processing"""
    try:
        print("🔄 Starting S3 background processing...")
        
        from src.agents.app_initialization import AppInitializationService
        
        init_service = AppInitializationService(agent.vector_store)
        
        if init_service.is_s3_enabled():
            print("🔄 S3 auto-loading enabled, initializing knowledge base...")
            
            init_results = init_service.initialize_knowledge_base(show_progress=False)
            
            if init_results['success']:
                callback('completed', init_results, None)
                print(f"✅ S3 knowledge base initialized: {init_results['message']}")
            else:
                callback('error', None, init_results['message'])
                print(f"❌ S3 initialization failed: {init_results['message']}")
        else:
            callback('disabled', None, None)
            print("ℹ️ S3 auto-loading not enabled")
            
    except Exception as e:
        callback('error', None, str(e))
        print(f"❌ S3 initialization error: {str(e)}")

def s3_callback(status, results, error):
    """Callback to update session state from background thread"""
    # Store in temporary files since we can't access session_state from thread
    import json
    with open('.s3_status.json', 'w') as f:
        json.dump({
            'status': status,
            'results': results,
            'error': error
        }, f, default=str)

def check_s3_background_status():
    """Check if S3 background processing has completed"""
    try:
        import json
        import os
        if os.path.exists('.s3_status.json'):
            with open('.s3_status.json', 'r') as f:
                data = json.load(f)
            
            # Update session state
            st.session_state.s3_status = data['status']
            if data['results']:
                st.session_state.s3_init_results = data['results']
            if data['error']:
                st.session_state.s3_init_error = data['error']
            
            # Clean up
            os.remove('.s3_status.json')
            return True
    except Exception as e:
        print(f"Error checking S3 status: {e}")
    return False

def start_s3_background_processing():
    """Start S3 processing in background after app is running"""
    if (st.session_state.agent and 
        not st.session_state.s3_background_started and
        st.session_state.s3_status == "ready"):
        
        st.session_state.s3_background_started = True
        st.session_state.s3_status = "processing"
        
        # Start background thread
        thread = threading.Thread(target=s3_background_thread, args=(st.session_state.agent, s3_callback), daemon=True)
        thread.start()
        print("🔄 S3 background processing thread started...")

def initialize_s3_background():
    """Legacy function - S3 no longer starts at initialization"""
    pass

def process_s3_if_queued():
    """Legacy function - replaced by fragment"""
    pass

def start_s3_after_ui_ready():
    """S3 processing is now manual only - no automatic startup processing"""
    pass

def main():
    """Main application function"""
    load_custom_css()
    initialize_session_state()
    
    # App starts immediately - NO S3 processing at startup
    
    # Create header - UI shows immediately
    create_header()
    
    # Create metrics dashboard
    create_metrics_dashboard()
    
    # Create main layout
    col1, col2 = st.columns([2, 1])
    
    with col1:
        create_chat_interface()
    
    with col2:
        create_sidebar()
    
    # Only start S3 processing after UI is fully rendered
    start_s3_after_ui_ready()
    
    # Footer
    st.markdown("---")
    st.markdown("""
    <div style="text-align: center; color: #6b7280; padding: 1rem;">
        <small>🏦 Financial Forecast AI | Built with Streamlit</small>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()