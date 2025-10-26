"""
S3 Document Loader for Financial Forecast AI
Handles automatic loading of documents from S3 bucket on application startup
"""

# Standard library imports
import hashlib
import io
import json
import mimetypes
import os
import re
from datetime import datetime
from typing import Dict, List, Optional, Tuple

# Third-party imports
import boto3
import markdown
import pandas as pd
import pypdf
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation


class S3DocumentLoader:
    """Loads documents from S3 bucket and processes them for the knowledge base"""
    
    def __init__(self, bucket_name: str = None, prefix: str = "documents/"):
        """
        Initialize S3 loader
        
        Args:
            bucket_name: S3 bucket name (from env if not provided)
            prefix: S3 prefix/folder to load documents from
        """
        self.bucket_name = bucket_name or os.getenv("S3_BUCKET_NAME")
        self.prefix = prefix
        self.s3_client = boto3.client('s3')
        
        if not self.bucket_name:
            raise ValueError("S3_BUCKET_NAME environment variable must be set")
    
    def list_s3_documents(self) -> List[Dict]:
        """List all documents in the S3 bucket with metadata"""
        try:
            response = self.s3_client.list_objects_v2(
                Bucket=self.bucket_name,
                Prefix=self.prefix
            )
            
            documents = []
            if 'Contents' in response:
                for obj in response['Contents']:
                    # Skip folders
                    if obj['Key'].endswith('/'):
                        continue
                    
                    # Get file extension and MIME type
                    file_extension = obj['Key'].split('.')[-1].lower()
                    mime_type, _ = mimetypes.guess_type(obj['Key'])
                    
                    documents.append({
                        'key': obj['Key'],
                        'size': obj['Size'],
                        'last_modified': obj['LastModified'],
                        'etag': obj['ETag'].strip('"'),
                        'file_extension': file_extension,
                        'mime_type': mime_type,
                        'filename': os.path.basename(obj['Key'])
                    })
            
            print(f"📁 Found {len(documents)} documents in S3 bucket: {self.bucket_name}")
            return documents
            
        except Exception as e:
            print(f"❌ Error listing S3 documents: {str(e)}")
            return []
    
    def download_document(self, s3_key: str) -> Tuple[bytes, Dict]:
        """Download a document from S3 and return content with metadata"""
        try:
            response = self.s3_client.get_object(Bucket=self.bucket_name, Key=s3_key)
            content = response['Body'].read()
            
            metadata = {
                'filename': os.path.basename(s3_key),
                's3_key': s3_key,
                's3_bucket': self.bucket_name,
                'size': len(content),
                'content_type': response.get('ContentType', ''),
                'last_modified': response.get('LastModified', ''),
                'etag': response.get('ETag', '').strip('"'),
                'source': 's3_auto_load'
            }
            
            return content, metadata
            
        except Exception as e:
            print(f"❌ Error downloading {s3_key}: {str(e)}")
            return b'', {}
    
    def extract_text_content(self, content: bytes, filename: str, mime_type: str = None) -> str:
        """Extract text content from various file formats"""
        file_extension = filename.split('.')[-1].lower()
        
        try:
            if file_extension == 'pdf' or (mime_type and 'pdf' in mime_type):
                return self._extract_pdf_text(content)
            
            elif file_extension in ['docx', 'doc'] or (mime_type and 'word' in mime_type):
                return self._extract_docx_text(content)
            
            elif file_extension in ['xlsx', 'xls'] or (mime_type and 'excel' in mime_type):
                return self._extract_excel_text(content)
            
            elif file_extension == 'pptx' or (mime_type and 'presentation' in mime_type):
                return self._extract_pptx_text(content)
            
            elif file_extension == 'csv':
                return self._extract_csv_text(content)
            
            elif file_extension == 'json':
                return self._extract_json_text(content)
            
            elif file_extension in ['html', 'htm']:
                return self._extract_html_text(content)
            
            elif file_extension == 'md':
                return self._extract_markdown_text(content)
            
            elif file_extension == 'rtf':
                return self._extract_rtf_text(content)
            
            else:
                # Try to decode as text
                return self._extract_plain_text(content)
                
        except Exception as e:
            return f"Error extracting text from {filename}: {str(e)}"
    
    def _extract_pdf_text(self, content: bytes) -> str:
        """Extract text from PDF content"""
        try:
            pdf_reader = pypdf.PdfReader(io.BytesIO(content))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            if not text.strip():
                return "PDF content could not be extracted (possibly scanned document)"
            return text
        except Exception as e:
            return f"Error processing PDF: {str(e)}"
    
    def _extract_docx_text(self, content: bytes) -> str:
        """Extract text from DOCX content"""
        try:
            doc = Document(io.BytesIO(content))
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract table content
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    text += row_text + "\n"
            
            return text if text.strip() else "No text content found in document"
        except Exception as e:
            return f"Error processing DOCX: {str(e)}"
    
    def _extract_excel_text(self, content: bytes) -> str:
        """Extract text from Excel content"""
        try:
            # Try different Excel formats
            try:
                df = pd.read_excel(io.BytesIO(content), sheet_name=None)
            except:
                df = pd.read_excel(io.BytesIO(content), engine='openpyxl', sheet_name=None)
            
            text = ""
            for sheet_name, sheet_df in df.items():
                text += f"Sheet: {sheet_name}\n"
                text += sheet_df.to_string(index=False) + "\n\n"
            
            return text
        except Exception as e:
            return f"Error processing Excel: {str(e)}"
    
    def _extract_pptx_text(self, content: bytes) -> str:
        """Extract text from PowerPoint content"""
        try:
            prs = Presentation(io.BytesIO(content))
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                text += "\n"
            
            return text
        except Exception as e:
            return f"Error processing PowerPoint: {str(e)}"
    
    def _extract_csv_text(self, content: bytes) -> str:
        """Extract text from CSV content"""
        try:
            df = pd.read_csv(io.BytesIO(content))
            return f"CSV Data ({len(df)} rows):\n" + df.to_string(index=False)
        except Exception as e:
            return f"Error processing CSV: {str(e)}"
    
    def _extract_json_text(self, content: bytes) -> str:
        """Extract text from JSON content"""
        try:
            json_data = json.loads(content.decode('utf-8'))
            return "JSON Content:\n" + json.dumps(json_data, indent=2)
        except Exception as e:
            return f"Error processing JSON: {str(e)}"
    
    def _extract_html_text(self, content: bytes) -> str:
        """Extract text from HTML content"""
        try:
            html_content = content.decode('utf-8')
            soup = BeautifulSoup(html_content, 'html.parser')
            return soup.get_text()
        except Exception as e:
            return f"Error processing HTML: {str(e)}"
    
    def _extract_markdown_text(self, content: bytes) -> str:
        """Extract text from Markdown content"""
        try:
            md_content = content.decode('utf-8')
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text()
        except Exception as e:
            return f"Error processing Markdown: {str(e)}"
    
    def _extract_rtf_text(self, content: bytes) -> str:
        """Extract text from RTF content"""
        try:
            rtf_content = content.decode('utf-8', errors='ignore')
            # Basic RTF text extraction (removes most RTF codes)
            text = re.sub(r'\\[a-z]+\d*\s?', ' ', rtf_content)
            text = re.sub(r'[{}]', '', text)
            text = re.sub(r'\s+', ' ', text).strip()
            return text
        except Exception as e:
            return f"Error processing RTF: {str(e)}"
    
    def _extract_plain_text(self, content: bytes) -> str:
        """Extract plain text content"""
        try:
            return content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                return content.decode('latin-1')
            except Exception as e:
                return f"Could not decode file content: {str(e)}"
    
    def generate_document_hash(self, content: str, metadata: Dict) -> str:
        """Generate a unique hash for the document to check for duplicates"""
        # Combine content and key metadata for hash
        hash_input = f"{content[:1000]}{metadata.get('filename', '')}{metadata.get('size', 0)}"
        return hashlib.md5(hash_input.encode()).hexdigest()
    
    def load_and_process_documents(self) -> List[Dict]:
        """
        Load all documents from S3, extract text, and prepare for indexing
        Returns list of processed documents with content and metadata
        """
        s3_documents = self.list_s3_documents()
        processed_documents = []
        
        if not s3_documents:
            print("📭 No documents found in S3 bucket")
            return []
        
        print(f"🔄 Processing {len(s3_documents)} documents from S3...")
        
        for doc_info in s3_documents:
            try:
                print(f"📄 Processing: {doc_info['filename']}")
                
                # Download document content
                content_bytes, metadata = self.download_document(doc_info['key'])
                
                if not content_bytes:
                    continue
                
                # Extract text content
                text_content = self.extract_text_content(
                    content_bytes, 
                    doc_info['filename'], 
                    doc_info['mime_type']
                )
                
                if not text_content or text_content.startswith("Error"):
                    print(f"⚠️ Failed to extract content from {doc_info['filename']}: {text_content}")
                    continue
                
                # Combine metadata
                full_metadata = {
                    **doc_info,
                    **metadata,
                    'content_length': len(text_content),
                    'processed_at': datetime.now().isoformat(),
                    'document_hash': self.generate_document_hash(text_content, metadata)
                }
                
                processed_documents.append({
                    'content': text_content,
                    'metadata': full_metadata
                })
                
                print(f"✅ Successfully processed: {doc_info['filename']} ({len(text_content)} chars)")
                
            except Exception as e:
                print(f"❌ Error processing {doc_info.get('filename', 'unknown')}: {str(e)}")
                continue
        
        print(f"🎉 Successfully processed {len(processed_documents)} documents from S3")
        return processed_documents