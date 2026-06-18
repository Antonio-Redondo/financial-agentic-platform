"""Path-based document loaders.

One entry point — :func:`load_text` — turns a file path into plain text.
Used by both the startup folder scanner and the Streamlit upload handler
(uploads are saved into ``DOCUMENTS_DIR`` first, then re-loaded from disk so
the two code paths stay identical).
"""
from __future__ import annotations

import io
import json
import os
import re
from typing import Callable, Dict


# Extensions we know how to extract. Anything else falls back to UTF-8 text.
SUPPORTED_EXTENSIONS = {
    "pdf", "docx", "doc", "xlsx", "xls", "csv", "pptx", "ppt",
    "txt", "md", "rtf", "html", "htm", "xml", "json",
}


def _read_bytes(path: str) -> bytes:
    with open(path, "rb") as f:
        return f.read()


def _load_pdf(path: str) -> str:
    import pypdf
    reader = pypdf.PdfReader(path)
    parts = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(parts).strip()


def _load_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" ".join(cell.text for cell in row.cells))
    return "\n".join(parts).strip()


def _load_excel(path: str) -> str:
    import pandas as pd
    sheets = pd.read_excel(path, sheet_name=None)
    parts = []
    for name, df in sheets.items():
        parts.append(f"Sheet: {name}")
        parts.append(df.to_string(index=False))
        parts.append("")
    return "\n".join(parts).strip()


def _load_csv(path: str) -> str:
    import pandas as pd
    df = pd.read_csv(path)
    return f"CSV Data ({len(df)} rows):\n{df.to_string(index=False)}"


def _load_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"Slide {i}:")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        parts.append("")
    return "\n".join(parts).strip()


def _load_json(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return "JSON Content:\n" + json.dumps(data, indent=2)


def _load_html(path: str) -> str:
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    return soup.get_text(separator="\n").strip()


def _load_markdown(path: str) -> str:
    import markdown
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        html = markdown.markdown(f.read())
    return BeautifulSoup(html, "html.parser").get_text(separator="\n").strip()


def _load_rtf(path: str) -> str:
    raw = _read_bytes(path).decode("utf-8", errors="ignore")
    text = re.sub(r"\\[a-z]+\d*\s?", " ", raw)
    text = re.sub(r"[{}]", "", text)
    return re.sub(r"\s+", " ", text).strip()


def _load_text(path: str) -> str:
    for enc in ("utf-8", "latin-1"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read().strip()
        except UnicodeDecodeError:
            continue
    return ""


_LOADERS: Dict[str, Callable[[str], str]] = {
    "pdf":  _load_pdf,
    "docx": _load_docx,
    "doc":  _load_docx,  # best-effort
    "xlsx": _load_excel,
    "xls":  _load_excel,
    "csv":  _load_csv,
    "pptx": _load_pptx,
    "ppt":  _load_pptx,
    "json": _load_json,
    "html": _load_html,
    "htm":  _load_html,
    "xml":  _load_html,
    "md":   _load_markdown,
    "rtf":  _load_rtf,
    "txt":  _load_text,
}


def load_text(path: str) -> str:
    """Extract text from ``path``. Unknown extensions are read as UTF-8 text.

    Returns an empty string if extraction fails; callers can then decide to
    skip the file rather than ingest a useless error message.
    """
    ext = os.path.splitext(path)[1].lstrip(".").lower()
    loader = _LOADERS.get(ext, _load_text)
    try:
        return loader(path)
    except Exception as e:  # noqa: BLE001 — surface the cause, never crash ingest
        print(f"⚠️  Loader failed for {path}: {e}", flush=True)
        return ""
